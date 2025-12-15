"""
Document processor for extracting and normalizing content from various formats
Enhanced with Azure AI Document Intelligence and Azure AI Vision
"""
import io
import logging
import base64
from typing import Dict, Any, Optional
from pathlib import Path

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import openpyxl

from azure.core.credentials import AzureKeyCredential
from azure.ai.documentintelligence import DocumentIntelligenceClient
from azure.ai.documentintelligence.models import AnalyzeDocumentRequest
from azure.ai.vision.imageanalysis import ImageAnalysisClient
from azure.ai.vision.imageanalysis.models import VisualFeatures

from src.models import DocumentType, ProcessedContent
from src.vision_analyzer import VisionAnalyzer

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Processes and extracts content from different document types"""
    
    def __init__(self, config=None):
        """
        Initialize document processor with optional Azure AI services
        
        Args:
            config: Config object with Azure AI service credentials
        """
        self.config = config
        self.doc_intelligence_client = None
        self.vision_client = None
        self.vision_analyzer = None
        
        # Initialize GPT-4 Vision analyzer if OpenAI is configured
        if config:
            try:
                self.vision_analyzer = VisionAnalyzer(config)
                logger.info("GPT-4 Vision analyzer enabled")
            except Exception as e:
                logger.warning(f"Failed to initialize GPT-4 Vision: {e}")
        
        # Initialize Azure Document Intelligence if configured
        if config and hasattr(config, 'azure_document_intelligence'):
            doc_config = config.azure_document_intelligence
            if doc_config.enabled and doc_config.endpoint and doc_config.api_key:
                try:
                    self.doc_intelligence_client = DocumentIntelligenceClient(
                        endpoint=doc_config.endpoint,
                        credential=AzureKeyCredential(doc_config.api_key)
                    )
                    logger.info("Azure Document Intelligence enabled")
                except Exception as e:
                    logger.warning(f"Failed to initialize Document Intelligence: {e}")
        
        # Initialize Azure AI Vision if configured
        if config and hasattr(config, 'azure_vision'):
            vision_config = config.azure_vision
            if vision_config.enabled and vision_config.endpoint and vision_config.api_key:
                try:
                    self.vision_client = ImageAnalysisClient(
                        endpoint=vision_config.endpoint,
                        credential=AzureKeyCredential(vision_config.api_key)
                    )
                    logger.info("Azure AI Vision enabled")
                except Exception as e:
                    logger.warning(f"Failed to initialize Azure AI Vision: {e}")
    
    def process(self, content: bytes, document_type: DocumentType, artifact_name: str) -> ProcessedContent:
        """
        Process document based on type
        
        Args:
            content: Raw document bytes
            document_type: Type of document
            artifact_name: Name of the artifact
            
        Returns:
            ProcessedContent with extracted information
        """
        processors = {
            DocumentType.PDF: self._process_pdf,
            DocumentType.DOCX: self._process_docx,
            DocumentType.PPTX: self._process_pptx,
            DocumentType.XLSX: self._process_xlsx,
            DocumentType.VSDX: self._process_vsdx,
            DocumentType.IMAGE: self._process_image,
            DocumentType.TEXT: self._process_text,
        }
        
        processor = processors.get(document_type, self._process_unknown)
        
        try:
            return processor(content, artifact_name)
        except Exception as e:
            logger.error(f"Error processing {artifact_name}: {e}")
            return ProcessedContent(
                artifact_name=artifact_name,
                document_type=document_type,
                extracted_text=f"Error processing document: {str(e)}",
                confidence_score=0.0
            )
    
    def _process_pdf(self, content: bytes, artifact_name: str) -> ProcessedContent:
        """Extract text from PDF using Document Intelligence or fallback to PyPDF2"""
        # Try Document Intelligence first for better quality
        if self.doc_intelligence_client:
            try:
                return self._process_with_document_intelligence(content, artifact_name, DocumentType.PDF)
            except Exception as e:
                logger.warning(f"Document Intelligence failed for {artifact_name}, using fallback: {e}")
        
        # Fallback to PyPDF2
        text_parts = []
        sections = []
        
        pdf_file = io.BytesIO(content)
        reader = PdfReader(pdf_file)
        
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            text_parts.append(page_text)
            sections.append({
                "type": "page",
                "number": i + 1,
                "content": page_text
            })
        
        full_text = "\n\n".join(text_parts)
        
        return ProcessedContent(
            artifact_name=artifact_name,
            document_type=DocumentType.PDF,
            extracted_text=full_text,
            sections=sections,
            structured_data={"page_count": len(reader.pages)},
            keywords=self._extract_keywords(full_text)
        )
    
    def _process_docx(self, content: bytes, artifact_name: str) -> ProcessedContent:
        """Extract text from DOCX"""
        doc_file = io.BytesIO(content)
        doc = Document(doc_file)
        
        paragraphs = []
        sections = []
        
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                paragraphs.append(para.text)
                sections.append({
                    "type": "paragraph",
                    "number": i + 1,
                    "content": para.text
                })
        
        full_text = "\n\n".join(paragraphs)
        
        # Extract tables if any
        tables_data = []
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                table_text.append(row_text)
            tables_data.append(table_text)
        
        return ProcessedContent(
            artifact_name=artifact_name,
            document_type=DocumentType.DOCX,
            extracted_text=full_text,
            sections=sections,
            structured_data={"paragraph_count": len(paragraphs), "table_count": len(tables_data)},
            keywords=self._extract_keywords(full_text)
        )
    
    def _process_pptx(self, content: bytes, artifact_name: str) -> ProcessedContent:
        """Extract text from PowerPoint"""
        ppt_file = io.BytesIO(content)
        prs = Presentation(ppt_file)
        
        text_parts = []
        sections = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            slide_content = "\n".join(slide_text)
            text_parts.append(slide_content)
            
            sections.append({
                "type": "slide",
                "number": i + 1,
                "content": slide_content
            })
        
        full_text = "\n\n".join(text_parts)
        
        return ProcessedContent(
            artifact_name=artifact_name,
            document_type=DocumentType.PPTX,
            extracted_text=full_text,
            sections=sections,
            structured_data={"slide_count": len(prs.slides)},
            keywords=self._extract_keywords(full_text)
        )
    
    def _process_image(self, content: bytes, artifact_name: str) -> ProcessedContent:
        """Process image using Azure AI Vision OCR or fallback to metadata"""
        # Try Azure AI Vision OCR first
        if self.vision_client:
            try:
                return self._process_with_vision_ocr(content, artifact_name)
            except Exception as e:
                logger.warning(f"Azure AI Vision failed for {artifact_name}, using fallback: {e}")
        
        # Fallback to basic image info
        image_file = io.BytesIO(content)
        img = Image.open(image_file)
        
        return ProcessedContent(
            artifact_name=artifact_name,
            document_type=DocumentType.IMAGE,
            extracted_text="[Image content - OCR not available. Enable Azure AI Vision for text extraction]",
            structured_data={
                "width": img.width,
                "height": img.height,
                "format": img.format
            },
            confidence_score=0.5
        )
    
    def _process_text(self, content: bytes, artifact_name: str) -> ProcessedContent:
        """Process plain text file"""
        text = content.decode('utf-8', errors='ignore')
        
        return ProcessedContent(
            artifact_name=artifact_name,
            document_type=DocumentType.TEXT,
            extracted_text=text,
            keywords=self._extract_keywords(text)
        )
    
    def _process_unknown(self, content: bytes, artifact_name: str) -> ProcessedContent:
        """Handle unknown document types"""
        return ProcessedContent(
            artifact_name=artifact_name,
            document_type=DocumentType.UNKNOWN,
            extracted_text="[Unknown document type]",
            confidence_score=0.0
        )
    
    def _process_xlsx(self, content: bytes, artifact_name: str) -> ProcessedContent:
        """Extract text and data from Excel files"""
        try:
            # Try Document Intelligence first for better table extraction
            if self.doc_intelligence_client:
                try:
                    return self._process_with_document_intelligence(content, artifact_name, DocumentType.XLSX)
                except Exception as e:
                    logger.warning(f"Document Intelligence failed for Excel, using openpyxl: {e}")
            
            # Fallback to openpyxl
            xlsx_file = io.BytesIO(content)
            workbook = openpyxl.load_workbook(xlsx_file, data_only=True)
            
            text_parts = []
            sections = []
            all_data = []
            
            for sheet in workbook.worksheets:
                sheet_data = []
                sheet_text = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_data):  # Skip empty rows
                        sheet_data.append(row_data)
                        sheet_text.append(" | ".join(row_data))
                
                if sheet_text:
                    sheet_content = "\n".join(sheet_text)
                    text_parts.append(f"Sheet: {sheet.title}\n{sheet_content}")
                    sections.append({
                        "type": "worksheet",
                        "name": sheet.title,
                        "content": sheet_content,
                        "row_count": len(sheet_data)
                    })
                    all_data.append({
                        "sheet": sheet.title,
                        "data": sheet_data
                    })
            
            full_text = "\n\n".join(text_parts)
            
            return ProcessedContent(
                artifact_name=artifact_name,
                document_type=DocumentType.XLSX,
                extracted_text=full_text,
                sections=sections,
                structured_data={
                    "sheet_count": len(workbook.worksheets),
                    "sheets": all_data
                },
                keywords=self._extract_keywords(full_text)
            )
        except Exception as e:
            logger.error(f"Error processing Excel file {artifact_name}: {e}")
            return ProcessedContent(
                artifact_name=artifact_name,
                document_type=DocumentType.XLSX,
                extracted_text=f"[Error processing Excel file: {str(e)}]",
                confidence_score=0.0
            )
    
    def _process_vsdx(self, content: bytes, artifact_name: str) -> ProcessedContent:
        """Process Visio diagrams (basic text extraction from XML)"""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            # VSDX is a ZIP file containing XML
            vsdx_file = io.BytesIO(content)
            text_parts = []
            
            with zipfile.ZipFile(vsdx_file, 'r') as zip_ref:
                # Extract text from pages
                for file_name in zip_ref.namelist():
                    if file_name.startswith('visio/pages/page') and file_name.endswith('.xml'):
                        xml_content = zip_ref.read(file_name)
                        try:
                            root = ET.fromstring(xml_content)
                            # Extract text from all text elements
                            for text_elem in root.iter():
                                if text_elem.text and text_elem.text.strip():
                                    text_parts.append(text_elem.text.strip())
                        except ET.ParseError:
                            continue
            
            full_text = "\n".join(text_parts) if text_parts else "[No text extracted from Visio diagram]"
            
            return ProcessedContent(
                artifact_name=artifact_name,
                document_type=DocumentType.VSDX,
                extracted_text=full_text,
                structured_data={"text_elements_found": len(text_parts)},
                keywords=self._extract_keywords(full_text),
                confidence_score=0.7
            )
        except Exception as e:
            logger.error(f"Error processing Visio file {artifact_name}: {e}")
            return ProcessedContent(
                artifact_name=artifact_name,
                document_type=DocumentType.VSDX,
                extracted_text=f"[Error processing Visio file: {str(e)}. Consider exporting to PDF or image format]",
                confidence_score=0.0
            )
    
    def _process_with_document_intelligence(self, content: bytes, artifact_name: str, 
                                           doc_type: DocumentType) -> ProcessedContent:
        """Process document using Azure Document Intelligence for superior extraction"""
        try:
            # Analyze document with layout model
            poller = self.doc_intelligence_client.begin_analyze_document(
                "prebuilt-layout",
                AnalyzeDocumentRequest(bytes_source=content)
            )
            result = poller.result()
            
            # Extract text with layout preserved
            full_text = result.content if hasattr(result, 'content') else ""
            
            # Extract tables
            tables_data = []
            if hasattr(result, 'tables'):
                for table in result.tables:
                    table_data = {
                        "row_count": table.row_count,
                        "column_count": table.column_count,
                        "cells": []
                    }
                    for cell in table.cells:
                        table_data["cells"].append({
                            "content": cell.content,
                            "row_index": cell.row_index,
                            "column_index": cell.column_index
                        })
                    tables_data.append(table_data)
            
            # Extract sections
            sections = []
            if hasattr(result, 'paragraphs'):
                for i, para in enumerate(result.paragraphs):
                    sections.append({
                        "type": "paragraph",
                        "number": i + 1,
                        "content": para.content,
                        "role": para.role if hasattr(para, 'role') else None
                    })
            
            return ProcessedContent(
                artifact_name=artifact_name,
                document_type=doc_type,
                extracted_text=full_text,
                sections=sections,
                structured_data={
                    "table_count": len(tables_data),
                    "tables": tables_data,
                    "page_count": len(result.pages) if hasattr(result, 'pages') else 0,
                    "processor": "Azure Document Intelligence"
                },
                keywords=self._extract_keywords(full_text),
                confidence_score=0.95
            )
        except Exception as e:
            logger.error(f"Document Intelligence processing failed: {e}")
            raise
    
    def _process_with_vision_ocr(self, content: bytes, artifact_name: str) -> ProcessedContent:
        """Process image using Azure AI Vision OCR and GPT-4 Vision analysis"""
        try:
            text_parts = []
            structured_data = {}
            
            # Step 1: Try Azure AI Vision OCR for text extraction
            if self.vision_client:
                result = self.vision_client.analyze(
                    image_data=content,
                    visual_features=[VisualFeatures.READ, VisualFeatures.CAPTION]
                )
                
                # Extract text from OCR
                if hasattr(result, 'read') and result.read:
                    for block in result.read.blocks:
                        for line in block.lines:
                            text_parts.append(line.text)
                
                # Get image description
                if hasattr(result, 'caption') and result.caption:
                    structured_data["caption"] = result.caption.text
            
            # Step 2: Use GPT-4 Vision for deeper analysis (diagrams, architecture)
            vision_analysis = ""
            if self.vision_analyzer:
                try:
                    # Determine analysis type based on file name
                    analysis_type = "general"
                    name_lower = artifact_name.lower()
                    if any(term in name_lower for term in ["architecture", "diagram", "design"]):
                        analysis_type = "architecture"
                    elif any(term in name_lower for term in ["network", "topology", "vnet"]):
                        analysis_type = "network"
                    elif any(term in name_lower for term in ["workflow", "process", "flow"]):
                        analysis_type = "workflow"
                    elif any(term in name_lower for term in ["whiteboard", "notes", "sketch"]):
                        # For whiteboards, extract text
                        vision_analysis = self.vision_analyzer.extract_text_from_whiteboard(content)
                    else:
                        # General diagram analysis
                        result = self.vision_analyzer.analyze_diagram(content, analysis_type)
                        if result.get("success"):
                            vision_analysis = result.get("analysis", "")
                            structured_data["vision_confidence"] = result.get("confidence")
                
                except Exception as e:
                    logger.warning(f"GPT-4 Vision analysis failed: {e}")
            
            # Combine OCR text and Vision analysis
            full_text_parts = []
            if text_parts:
                full_text_parts.append("=== Extracted Text ===\n" + "\n".join(text_parts))
            if vision_analysis:
                full_text_parts.append("=== Visual Analysis ===\n" + vision_analysis)
            
            full_text = "\n\n".join(full_text_parts) if full_text_parts else "[No text detected in image]"
            
            # Get image metadata
            image_file = io.BytesIO(content)
            img = Image.open(image_file)
            
            structured_data.update({
                "width": img.width,
                "height": img.height,
                "format": img.format,
                "text_blocks_found": len(text_parts),
                "processor": "Azure AI Vision + GPT-4 Vision" if vision_analysis else "Azure AI Vision"
            })
            
            return ProcessedContent(
                artifact_name=artifact_name,
                document_type=DocumentType.IMAGE,
                extracted_text=full_text,
                structured_data=structured_data,
                keywords=self._extract_keywords(full_text),
                confidence_score=0.95 if vision_analysis else 0.9
            )
        except Exception as e:
            logger.error(f"Image processing failed: {e}")
            raise
    
    def _extract_keywords(self, text: str) -> list:
        """Simple keyword extraction (can be enhanced with NLP)"""
        # Basic implementation - can be improved with Azure AI Language
        keywords = []
        
        # Azure-related keywords
        azure_keywords = [
            "subscription", "resource group", "virtual network", "vnet",
            "subnet", "nsg", "security", "compliance", "governance",
            "landing zone", "management group", "policy", "rbac",
            "storage account", "key vault", "application gateway",
            "load balancer", "firewall", "vpn", "expressroute"
        ]
        
        text_lower = text.lower()
        for keyword in azure_keywords:
            if keyword in text_lower:
                keywords.append(keyword)
        
        return list(set(keywords))
